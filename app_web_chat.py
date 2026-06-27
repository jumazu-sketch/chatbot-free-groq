import base64
import json
import os
from pathlib import Path
from urllib import error, request

import streamlit as st

GROQ_API_URL = "https://api.groq.com/openai/v1/chat/completions"
GROQ_KEYS_URL = "https://console.groq.com/keys"
DEFAULT_MODEL = "meta-llama/llama-4-scout-17b-16e-instruct"
DEFAULT_API_KEY = "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"
AVAILABLE_MODELS = [
    "meta-llama/llama-4-scout-17b-16e-instruct",
    "llama-3.2-11b-vision-preview",
    "llama-3.2-90b-vision-preview",
    "llama-3.1-8b-instant",
    "llama-3.3-70b-versatile",
    "groq/compound",
    "groq/compound-mini",
    "openai/gpt-oss-120b",
    "openai/gpt-oss-20b",
    "qwen/qwen3-32b",
    "whisper-large-v3",
    "whisper-large-v3-turbo",
]

# Banco de dados de Assistentes / Personas disponíveis no chat
ASSISTANTS = {
    "Assistente Padrão": {
        "description": "Um assistente geral, objetivo e didático. Responde de forma curta e clara em português.",
        "prompt": (
            "Voce e um assistente objetivo e didatico. "
            "Responda em portugues de forma clara e curta quando possivel."
        ),
        "starters": [
            {"label": "📝 Melhore uma frase", "text": "Melhore esta frase em Inglês para soar mais natural e profissional: "},
            {"label": "🗣️ Tradução perfeita", "text": "Qual é a melhor tradução para o português da frase/termo: "},
            {"label": "📄 Resumir texto", "text": "Resuma o seguinte texto de forma didática e em tópicos: "}
        ]
    },
    "Natural Clear English": {
        "description": "Designed to bridge the gap between non-native English and professional clarity by correcting verb tenses and idioms while preserving simplicity.",
        "prompt": (
            "You are a \"Natural English Editor\" for a non-native professional. "
            "Your goal is to help the user polish their emails and messages so they are grammatically correct and clear, "
            "while keeping the tone simple, direct, and authentic to a non-native speaker.\n\n"
            "Objectives:\n"
            "1. Fix Essential Errors: Correct \"really needed\" mistakes like verb tenses (e.g., change \"I'd like to aligned\" to \"I'd like to align\") and incorrect prepositions.\n"
            "2. Use Common Expressions: Replace awkward phrasing with common, everyday workplace idioms (e.g., change \"in this meanwhile\" to \"in the meantime\").\n"
            "3. Keep it Simple: Do NOT use \"big\" or sophisticated words. Avoid sounding like a corporate lawyer or a professional writer. Use high-frequency vocabulary.\n"
            "4. Preserve the User's Voice: Keep the message short and to the point. The final result should sound like a competent professional who is a non-native speaker, not a native poet.\n\n"
            "Output Format:\n"
            "- First, provide the \"Suggested Version\" (the corrected text).\n"
            "- Second, provide a \"Why?\" section with a few bullet points explaining the most important corrections in simple terms."
        ),
        "starters": [
            {"label": "✉️ Polish an email", "text": "Please polish this email for me to make it grammatically correct and clear:\n\n"},
            {"label": "💬 Fix this sentence", "text": "Fix this sentence to sound simple and correct:\n\n"},
            {"label": "🤝 Business request", "text": "Make this business request clear, direct, and natural:\n\n"}
        ]
    }
}

APP_SETTINGS_DIR = Path(os.getenv("APPDATA", Path.home())) / "TesteSolver"
APP_SETTINGS_FILE = APP_SETTINGS_DIR / "config.json"


def load_settings() -> dict:
    if not APP_SETTINGS_FILE.exists():
        return {}

    try:
        return json.loads(APP_SETTINGS_FILE.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return {}


def save_settings(settings: dict) -> None:
    APP_SETTINGS_DIR.mkdir(parents=True, exist_ok=True)
    APP_SETTINGS_FILE.write_text(
        json.dumps(settings, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


def validate_api_key_format(api_key: str) -> bool:
    return bool(api_key) and api_key.startswith("gsk_") and len(api_key) >= 20


def extract_text_from_content(content):
    if isinstance(content, str):
        return content.strip()

    if isinstance(content, list):
        parts = []
        for part in content:
            if isinstance(part, dict):
                text = part.get("text")
                if isinstance(text, str):
                    parts.append(text)
        return "".join(parts).strip()

    return ""


def extract_text_from_file(uploaded_file) -> str:
    filename = uploaded_file.name.lower()
    
    # 1. Arquivos de Texto (.txt)
    if filename.endswith(".txt"):
        return uploaded_file.read().decode("utf-8", errors="replace")
        
    # 2. Arquivos PDF (.pdf)
    elif filename.endswith(".pdf"):
        try:
            import pypdf
            pdf_reader = pypdf.PdfReader(uploaded_file)
            text = ""
            for page in pdf_reader.pages:
                text_page = page.extract_text()
                if text_page:
                    text += text_page + "\n"
            return text
        except ImportError:
            return "Erro: Para ler arquivos PDF, a biblioteca 'pypdf' precisa estar instalada no requirements.txt."
        except Exception as e:
            return f"Erro ao ler PDF: {str(e)}"
            
    # 3. Arquivos CSV (.csv)
    elif filename.endswith(".csv"):
        try:
            import pandas as pd
            df = pd.read_csv(uploaded_file)
            return df.to_string(index=False)
        except ImportError:
            return "Erro: Para ler arquivos CSV, a biblioteca 'pandas' precisa estar instalada no requirements.txt."
        except Exception as e:
            return f"Erro ao ler CSV: {str(e)}"
            
    # 4. Arquivos Excel (.xlsx, .xls)
    elif filename.endswith((".xlsx", ".xls")):
        if filename.endswith(".xls"):
            return "Erro: O formato antigo '.xls' não é suportado diretamente. Por favor, salve seu arquivo como '.xlsx' (Excel Moderno) e envie novamente."
        try:
            import pandas as pd
            df = pd.read_excel(uploaded_file)
            return df.to_string(index=False)
        except ImportError:
            return "Erro: Para ler arquivos Excel, as bibliotecas 'pandas' e 'openpyxl' precisam estar instaladas no requirements.txt."
        except Exception as e:
            return f"Erro ao ler Excel: {str(e)}"
            
    # 5. Arquivos Word (.docx, .doc)
    elif filename.endswith((".docx", ".doc")):
        if filename.endswith(".doc"):
            return "Erro: O formato antigo '.doc' não é suportado diretamente. Por favor, salve seu arquivo como '.docx' (Word Moderno) e envie novamente."
        try:
            import docx
            doc = docx.Document(uploaded_file)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            return "\n".join(full_text)
        except ImportError:
            return "Erro: Para ler arquivos Word, a biblioteca 'python-docx' precisa estar instalada no requirements.txt."
        except Exception as e:
            return f"Erro ao ler DOCX: {str(e)}"
            
    # 6. Arquivos PowerPoint (.pptx, .ppt)
    elif filename.endswith((".pptx", ".ppt")):
        if filename.endswith(".ppt"):
            return "Erro: O formato antigo '.ppt' não é suportado diretamente. Por favor, salve seu arquivo como '.pptx' (PowerPoint Moderno) e envie novamente."
        try:
            from pptx import Presentation
            prs = Presentation(uploaded_file)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        except ImportError:
            return "Erro: Para ler arquivos PowerPoint, a biblioteca 'python-pptx' precisa estar instalada no requirements.txt."
        except Exception as e:
            return f"Erro ao ler PPTX: {str(e)}"
            
    return ""


def ask_groq(api_key: str, model: str, messages: list[dict]) -> str:
    payload = {
        "model": model,
        "messages": messages,
        "temperature": 0.2,
        "max_completion_tokens": 900,
        "top_p": 1,
        "stream": False,
    }

    api_request = request.Request(
        GROQ_API_URL,
        data=json.dumps(payload).encode("utf-8"),
        headers={
            "Authorization": f"Bearer {api_key}",
            "Accept": "application/json",
            "Content-Type": "application/json",
            "User-Agent": "Mozilla/5.0",
        },
        method="POST",
    )

    try:
        with request.urlopen(api_request, timeout=60) as response:
            data = json.loads(response.read().decode("utf-8"))
    except error.HTTPError as exc:
        details = exc.read().decode("utf-8", errors="replace")
        try:
            parsed = json.loads(details)
            details = parsed.get("error", {}).get("message", details)
        except json.JSONDecodeError:
            pass
        raise RuntimeError(f"Groq retornou erro HTTP {exc.code}: {details}") from exc
    except error.URLError as exc:
        raise RuntimeError(f"Falha de conexao com a API Groq: {exc.reason}") from exc

    choices = data.get("choices") or []
    if not choices:
        raise RuntimeError("A API Groq nao retornou resposta.")

    content = choices[0].get("message", {}).get("content", "")
    answer = extract_text_from_content(content)
    if not answer:
        raise RuntimeError("A API Groq respondeu sem conteudo.")

    return answer


def build_messages(history: list[dict], system_prompt: str) -> list[dict]:
    api_messages = [{"role": "system", "content": system_prompt}]
    for msg in history:
        api_messages.append({
            "role": msg["role"],
            "content": msg["content"]
        })
    return api_messages


def init_state() -> None:
    if "settings" not in st.session_state:
        st.session_state.settings = load_settings()

    if "messages" not in st.session_state:
        st.session_state.messages = []


def get_effective_key() -> str:
    typed = st.session_state.get("typed_key", "").strip()
    if typed:
        return typed

    saved = st.session_state.settings.get("api_key", "")
    if isinstance(saved, str) and saved.strip():
        return saved.strip()

    return DEFAULT_API_KEY


def sidebar_settings() -> tuple[str, str, str]:
    st.sidebar.title("Configuracao")
    st.sidebar.markdown(f"Criar chave: {GROQ_KEYS_URL}")

    saved_key = st.session_state.settings.get("api_key", "")
    if not isinstance(saved_key, str) or not saved_key:
        saved_key = DEFAULT_API_KEY

    typed_key = st.sidebar.text_input(
        "API Key da Groq",
        value=saved_key,
        type="password",
        key="typed_key",
        placeholder="gsk_...",
    )

    # Seletor de Modelo
    def format_model_name(model_id: str) -> str:
        if "vision" in model_id.lower() or "scout" in model_id.lower():
            return f"📷 [MULTIMODAL] {model_id}"
        return f"📄 [TEXTO] {model_id}"

    model = st.sidebar.selectbox(
        "Modelo",
        options=AVAILABLE_MODELS,
        index=AVAILABLE_MODELS.index(DEFAULT_MODEL),
        format_func=format_model_name,
    )

    st.sidebar.markdown("---")
    st.sidebar.subheader("Persona do Chatbot")
    
    # Seletor do Assistente / Persona
    selected_assistant = st.sidebar.selectbox(
        "Selecione o Assistente",
        options=list(ASSISTANTS.keys()),
        index=0
    )
    
    # Mostra a descrição do assistente selecionado na barra lateral
    st.sidebar.info(ASSISTANTS[selected_assistant]["description"])

    st.sidebar.markdown("---")

    if st.sidebar.button("Salvar chave"):
        if not validate_api_key_format(typed_key.strip()):
            st.sidebar.error("Chave invalida. Ela normalmente comeca com gsk_.")
        else:
            st.session_state.settings["api_key"] = typed_key.strip()
            try:
                save_settings(st.session_state.settings)
                st.sidebar.success("Chave salva no AppData.")
            except OSError as exc:
                st.sidebar.error(f"Nao foi possivel salvar: {exc}")

    if st.sidebar.button("Limpar conversa"):
        st.session_state.messages = []
        st.rerun()

    return typed_key.strip(), model, selected_assistant


def render_chat() -> None:
    for msg in st.session_state.messages:
        role = msg.get("role", "assistant")
        content = msg.get("content", "")
        display = msg.get("display_content", "")
        with st.chat_message(role):
            if role == "user":
                if isinstance(content, list):
                    if display:
                        st.markdown(display)
                    else:
                        for item in content:
                            if item.get("type") == "text":
                                st.markdown(item.get("text"))
                    for item in content:
                        if item.get("type") == "image_url":
                            st.image(item.get("image_url", {}).get("url", ""), width=300)
                else:
                    st.markdown(display or content)
            else:
                st.markdown(content)


def main() -> None:
    st.set_page_config(page_title="Mazzucas's LLM", page_icon=":speech_balloon:", layout="centered")
    init_state()

    st.title("Mazzuca´s bot")
    st.caption("Aplicacao web minima de chat usando sua chave da Groq.")

    typed_key, selected_model, selected_assistant = sidebar_settings()

    env_key = os.getenv("GROQ_API_KEY", "").strip()
    effective_key = typed_key or get_effective_key() or env_key

    if not validate_api_key_format(effective_key):
        st.info("Informe e salve uma API key valida para comecar o chat.")

    render_chat()

    # Bloco de Sugestões de Conversa Dinâmico (muda dependendo do Assistente ativo)
    if not st.session_state.messages:
        st.markdown("### 💡 Sugestões de conversa:")
        starters = ASSISTANTS[selected_assistant]["starters"]
        col1, col2, col3 = st.columns(3)
        
        with col1:
            if st.button(starters[0]["label"], use_container_width=True):
                st.session_state.chat_input = starters[0]["text"]
                st.rerun()
        with col2:
            if st.button(starters[1]["label"], use_container_width=True):
                st.session_state.chat_input = starters[1]["text"]
                st.rerun()
        with col3:
            if st.button(starters[2]["label"], use_container_width=True):
                st.session_state.chat_input = starters[2]["text"]
                st.rerun()

    # O chat_input com a chave (key) para permitir o preenchimento automático pelas sugestões
    prompt = st.chat_input(
        "Digite sua pergunta ou envie arquivos...",
        accept_file="multiple",
        file_type=["txt", "pdf", "csv", "xlsx", "xls", "docx", "doc", "pptx", "ppt", "png", "jpg", "jpeg"],
        key="chat_input"
    )
    
    if not prompt:
        return

    user_text = (prompt.text or "").strip()
    uploaded_files = prompt.files or []

    if not user_text and uploaded_files:
        user_text = "Analise o arquivo que enviei e resuma suas principais informações."

    file_context_list = []
    image_base64_list = []
    processed_file_names = []
    
    for uploaded_file in uploaded_files:
        file_name = uploaded_file.name.lower()
        if file_name.endswith((".png", ".jpg", ".jpeg")):
            bytes_data = uploaded_file.read()
            base64_str = base64.b64encode(bytes_data).decode("utf-8")
            ext = "png" if file_name.endswith(".png") else "jpeg"
            image_base64_list.append(f"data:image/{ext};base64,{base64_str}")
            processed_file_names.append(uploaded_file.name)
        else:
            text_extracted = extract_text_from_file(uploaded_file)
            if text_extracted and not text_extracted.startswith("Erro:"):
                file_context_list.append(f"--- Conteúdo do arquivo '{uploaded_file.name}' ---\n{text_extracted}\n---")
                processed_file_names.append(uploaded_file.name)
            elif text_extracted.startswith("Erro:"):
                st.error(text_extracted)
                return

    file_context = "\n\n".join(file_context_list)
    
    if image_base64_list:
        is_vision = "vision" in selected_model or "scout" in selected_model
        if not is_vision:
            st.warning("⚠️ O modelo atual pode não aceitar imagens. Escolha um modelo de visão no menu lateral.")
            
        final_prompt = user_text
        if file_context:
            final_prompt = f"{file_context}\n\nPergunta do usuário: {user_text}"
            
        message_content = [{"type": "text", "text": final_prompt}]
        for img in image_base64_list:
            message_content.append({"type": "image_url", "image_url": {"url": img}})
            
        files_str = ", ".join(processed_file_names)
        display_content = f"📎 Arquivos enviados: *{files_str}*\n\n{user_text}"
    elif file_context:
        message_content = f"{file_context}\n\nPergunta do usuário: {user_text}"
        files_str = ", ".join(processed_file_names)
        display_content = f"📄 Arquivos enviados: *{files_str}*\n\n{user_text}"
    else:
        message_content = user_text
        display_content = user_text

    st.session_state.messages.append({
        "role": "user",
        "content": message_content,
        "display_content": display_content
    })
    
    with st.chat_message("user"):
        if isinstance(message_content, list):
            if display_content:
                st.markdown(display_content)
            for item in message_content:
                if item["type"] == "image_url":
                    st.image(item["image_url"]["url"], width=300)
        else:
            st.markdown(display_content)

    if not validate_api_key_format(effective_key):
        with st.chat_message("assistant"):
            st.error("API key invalida. Salve uma chave da Groq para continuar.")
        st.session_state.messages.append(
            {"role": "assistant", "content": "API key invalida. Configure uma chave valida."}
        )
        return

    with st.chat_message("assistant"):
        with st.spinner("Pensando..."):
            try:
                # Carrega o prompt do assistente dinamicamente baseado na seleção da sidebar
                active_system_prompt = ASSISTANTS[selected_assistant]["prompt"]
                api_messages = build_messages(st.session_state.messages, active_system_prompt)
                answer = ask_groq(effective_key, selected_model, api_messages)
            except Exception as exc:
                st.error(str(exc))
                st.session_state.messages.append(
                    {"role": "assistant", "content": f"Erro: {exc}"}
                )
                return

        st.markdown(answer)

    st.session_state.messages.append({"role": "assistant", "content": answer})


if __name__ == "__main__":
    main()
